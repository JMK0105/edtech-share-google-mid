from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Pt

def replace_text_preserve_style(text_frame, new_text):
    DEFAULT_FONT_NAME = "맑은 고딕"
    DEFAULT_FONT_SIZE = Pt(16)
    DEFAULT_FONT_COLOR = RGBColor(0, 0, 0)

    # 줄바꿈 문자열 처리
    lines = new_text.replace("\n", "\\n").split("\\n")
    lines = [line.strip() for line in lines if line.strip()]
    if not lines:
        lines = ["(내용 없음)"]

    # 기존 스타일 추출
    template_run = None
    template_align = None
    if text_frame.paragraphs and text_frame.paragraphs[0].runs:
        template_run = text_frame.paragraphs[0].runs[0]
        template_align = text_frame.paragraphs[0].alignment

    # 기존 텍스트 제거 (빈 문단 하나는 자동 생성됨)
    text_frame.clear()
    first_paragraph = text_frame.paragraphs[0]

    for i, line in enumerate(lines):
        # 첫 줄은 기존 paragraph 사용
        p = first_paragraph if i == 0 else text_frame.add_paragraph()
        if template_align is not None:
            p.alignment = template_align
        run = p.add_run()
        run.text = line

        # 스타일 복사 또는 fallback
        if template_run:
            run.font.name = template_run.font.name
            run.font.size = template_run.font.size or DEFAULT_FONT_SIZE
            run.font.bold = template_run.font.bold
            run.font.italic = template_run.font.italic
            try:
                run.font.color.rgb = template_run.font.color.rgb or DEFAULT_FONT_COLOR
            except:
                run.font.color.rgb = DEFAULT_FONT_COLOR
        else:
            run.font.name = DEFAULT_FONT_NAME
            run.font.size = DEFAULT_FONT_SIZE
            run.font.color.rgb = DEFAULT_FONT_COLOR

        if i > 0:
            p.space_before = Pt(6)

    text_frame.word_wrap = True
    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

def insert_structured_content(template_path, structured_slides):
    """
    템플릿 PPT 파일을 열고, structured_slides 데이터를 슬라이드에 삽입합니다.
    """
    prs = Presentation(template_path)

    for i, slide_data in enumerate(structured_slides):
        if i >= len(prs.slides):
            break

        slide = prs.slides[i]
        title_kr = slide_data.get("title_kr", "")
        title_en = slide_data.get("title_en", "")
        content = slide_data.get("content", "")
        keywords = slide_data.get("keywords", "")

        # 디버깅용 로그 (선택적)
        if not any([title_kr, title_en, keywords, content]):
            print(f"⚠️ 슬라이드 {i+1}에 입력된 값이 없습니다.")

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            if shape.name == "TitleKRBox":
                replace_text_preserve_style(shape.text_frame, title_kr)
            elif shape.name == "TitleENBox":
                replace_text_preserve_style(shape.text_frame, title_en)
            elif shape.name == "BodyBox":
                replace_text_preserve_style(shape.text_frame, content)
            elif shape.name == "KeywordBox":
                replace_text_preserve_style(shape.text_frame, keywords)

    return prs
